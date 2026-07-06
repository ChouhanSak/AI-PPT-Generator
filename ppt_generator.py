"""
ppt_generator.py
Premium AI presentation engine.
Gemini creates structured slide strategy; python-pptx renders varied layouts.
"""

import json
import os
import re
from io import BytesIO

from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

MODEL = "gemini-3.5-flash"

SYSTEM_PROMPT = """You are a subject-matter expert and senior presentation designer.
Create a concise, accurate, visually strategic presentation outline.
Each slide must communicate one clear idea. Prefer concrete mechanisms, categories,
stages, comparisons, examples, and genuine statistics. Never use generic filler.
Choose layouts based on communication purpose, not randomly.
Return ONLY valid JSON matching the requested schema."""

def _build_user_prompt(title: str, num_slides: int, audience: str, tone: str) -> str:
    return f"""
Presentation title: "{title}"
Target audience: {audience}
Tone: {tone}
Number of CONTENT slides: {num_slides}

Return exactly this JSON structure:
{{
  "palette": {{
    "primary": "#RRGGBB",
    "secondary": "#RRGGBB",
    "accent": "#RRGGBB",
    "surface": "#RRGGBB",
    "text": "#RRGGBB"
  }},
  "subtitle": "short title-slide subtitle",
  "slides": [
    {{
      "layout": "hero_statement|bullets|cards|two_column|stat|timeline|process|problem_solution|three_metrics|quote",
      "eyebrow": "2-4 word section label",
      "heading": "short strong slide headline",
      "subheading": "optional one-line context",
      "bullets": ["specific point", "specific point", "specific point"],
      "items": [
        {{"title": "short title", "body": "short specific explanation"}}
      ],
      "stat": "73%",
      "stat_label": "short factual label",
      "left_title": "left concept",
      "left_bullets": ["specific point", "specific point", "specific point"],
      "right_title": "right concept",
      "right_bullets": ["specific point", "specific point", "specific point"],
      "quote": "short quotation only when a reliable attribution is known",
      "quote_source": "source/person",
      "metrics": [
        {{"value": "42%", "label": "short label"}}
      ]
    }}
  ],
  "conclusion": {{
    "heading": "short closing headline",
    "bullets": ["takeaway", "takeaway", "takeaway"]
  }}
}}

Design/content rules:
- Cover the topic's real sub-parts across exactly {num_slides} content slides.
- Choose a sophisticated topic-aware palette with strong contrast.
- Vary layouts. Do not use the same layout more than twice consecutively.
- hero_statement: one major insight with 2-4 supporting bullets.
- cards: 3 or 4 categories; use items.
- two_column: a genuine contrast, pair, before/after, or two-sided concept.
- stat: only for a genuine defensible number; use stat, stat_label, bullets.
- timeline: chronological stages; use 3-5 items.
- process: ordered workflow; use 3-5 items.
- problem_solution: left side problem, right side solution.
- three_metrics: exactly 3 genuine metrics.
- quote: only if attribution is reliable; never invent quotations.
- bullets: 3-5 concise specific points.
- Keep body copy concise. Avoid paragraphs and filler.
- Output pure JSON only.
"""

def _extract_json(text: str) -> dict:
    cleaned = re.sub(r"^```(?:json)?|```$", "", text.strip(), flags=re.MULTILINE).strip()
    return json.loads(cleaned)

def generate_outline(title: str, num_slides: int = 6,
                     audience: str = "general professional audience",
                     tone: str = "confident and clear") -> dict:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is missing from .env file.")

    client = genai.Client(api_key=api_key)
    try:
        print(f"Attempting content generation using: {MODEL}...")
        response = client.models.generate_content(
            model=MODEL,
            contents=_build_user_prompt(title, num_slides, audience, tone),
            config=types.GenerateContentConfig(
                system_instruction=SYSTEM_PROMPT,
                response_mime_type="application/json",
                temperature=0.7,
                max_output_tokens=5000,
            ),
        )
        if not response.text:
            raise RuntimeError("Gemini returned an empty response.")
        outline = _extract_json(response.text)
        if not isinstance(outline.get("slides"), list):
            raise RuntimeError("Gemini returned an invalid slide structure.")
        return outline
    except Exception as e:
        print("\n========== GEMINI REAL ERROR ==========")
        print(type(e).__name__)
        print(str(e))
        print("=======================================\n")
        raise RuntimeError(f"Gemini generation failed: {e}") from e


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_FONT = "Aptos Display"
BODY_FONT = "Aptos"

WHITE = RGBColor(255, 255, 255)
INK = RGBColor(28, 31, 38)
MUTED = RGBColor(101, 107, 120)
LIGHT = RGBColor(246, 247, 249)

def _hex_to_rgb(value: str, fallback="#1E2761") -> RGBColor:
    try:
        value = (value or fallback).lstrip("#")
        if len(value) != 6:
            raise ValueError
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except (ValueError, TypeError):
        value = fallback.lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

def _set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _shape(slide, shape_type, left, top, width, height, fill, radius_line=None):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def _text(slide, left, top, width, height, text, size, color,
          bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _bullets(slide, left, top, width, height, bullets, color=INK, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate((bullets or [])[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(13)
    return box

def _eyebrow(slide, data, primary):
    label = str(data.get("eyebrow", "INSIGHT")).upper()
    _text(slide, Inches(.72), Inches(.42), Inches(3.5), Inches(.3),
          label, 10, primary, bold=True)

def _heading(slide, data, primary):
    _eyebrow(slide, data, primary)
    _text(slide, Inches(.7), Inches(.82), Inches(11.8), Inches(.8),
          data.get("heading", ""), 29, INK, bold=True, font=HEADER_FONT)

def _footer(slide, number, primary):
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(.7), Inches(7.12),
           Inches(10.8), Inches(.025), primary)
    _text(slide, Inches(11.7), Inches(6.96), Inches(.8), Inches(.3),
          f"{number:02d}", 10, MUTED, align=PP_ALIGN.RIGHT)

def _title_slide(prs, title, subtitle, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, primary)
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(9.8), Inches(0),
           Inches(3.533), Inches(7.5), secondary)
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(.85), Inches(1.25),
           Inches(.08), Inches(4.6), accent)
    _text(slide, Inches(1.25), Inches(1.65), Inches(7.8), Inches(2.3),
          title, 42, WHITE, bold=True, font=HEADER_FONT)
    _text(slide, Inches(1.27), Inches(4.25), Inches(7.5), Inches(.8),
          subtitle, 17, WHITE)
    _text(slide, Inches(1.27), Inches(6.45), Inches(4), Inches(.35),
          "AI-GENERATED PRESENTATION", 9, accent, bold=True)
    return slide

def _hero_statement_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, primary)
    _text(slide, Inches(.8), Inches(.55), Inches(3), Inches(.35),
          str(data.get("eyebrow", "KEY IDEA")).upper(), 10, accent, bold=True)
    _text(slide, Inches(.8), Inches(1.25), Inches(10.8), Inches(2.2),
          data.get("heading", ""), 38, WHITE, bold=True, font=HEADER_FONT)
    _text(slide, Inches(.82), Inches(3.65), Inches(9.8), Inches(.65),
          data.get("subheading", ""), 16, WHITE)
    _bullets(slide, Inches(7.4), Inches(4.65), Inches(5), Inches(1.8),
             data.get("bullets", []), WHITE, 15)
    return slide

def _bullets_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _heading(slide, data, primary)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.7), Inches(1.9),
           Inches(3.1), Inches(4.6), primary)
    _text(slide, Inches(1.05), Inches(2.35), Inches(2.4), Inches(1.5),
          data.get("subheading") or data.get("heading", ""), 24, WHITE,
          bold=True, font=HEADER_FONT)
    _bullets(slide, Inches(4.45), Inches(2.05), Inches(7.7), Inches(4.5),
             data.get("bullets", []), INK, 18)
    return slide

def _cards_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, LIGHT)
    _heading(slide, data, primary)
    items = (data.get("items") or [])[:4]
    count = max(1, len(items))
    gap = .22
    card_w = (11.93 - gap * (count - 1)) / count
    for i, item in enumerate(items):
        x = .7 + i * (card_w + gap)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0),
               Inches(card_w), Inches(4.45), WHITE)
        _text(slide, Inches(x + .28), Inches(2.35), Inches(.6), Inches(.5),
              f"{i+1:02d}", 12, primary, bold=True)
        _text(slide, Inches(x + .28), Inches(3.0), Inches(card_w-.56), Inches(.9),
              item.get("title", ""), 20, INK, bold=True, font=HEADER_FONT)
        _text(slide, Inches(x + .28), Inches(4.15), Inches(card_w-.56), Inches(1.55),
              item.get("body", ""), 14, MUTED)
    return slide

def _two_column_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _heading(slide, data, primary)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.7), Inches(1.9),
           Inches(5.75), Inches(4.7), LIGHT)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.9),
           Inches(5.85), Inches(4.7), primary)
    _text(slide, Inches(1.05), Inches(2.25), Inches(4.9), Inches(.6),
          data.get("left_title", ""), 22, INK, bold=True, font=HEADER_FONT)
    _bullets(slide, Inches(1.05), Inches(3.05), Inches(4.9), Inches(3),
             data.get("left_bullets", []), INK, 15)
    _text(slide, Inches(7.1), Inches(2.25), Inches(4.9), Inches(.6),
          data.get("right_title", ""), 22, WHITE, bold=True, font=HEADER_FONT)
    _bullets(slide, Inches(7.1), Inches(3.05), Inches(4.9), Inches(3),
             data.get("right_bullets", []), WHITE, 15)
    return slide

def _stat_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _heading(slide, data, primary)
    _text(slide, Inches(.75), Inches(2.0), Inches(5.2), Inches(1.6),
          data.get("stat", ""), 64, primary, bold=True, font=HEADER_FONT)
    _text(slide, Inches(.8), Inches(3.55), Inches(4.7), Inches(1),
          data.get("stat_label", ""), 16, MUTED)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.95),
           Inches(6.1), Inches(4.5), LIGHT)
    _bullets(slide, Inches(6.45), Inches(2.55), Inches(5.15), Inches(3.4),
             data.get("bullets", []), INK, 17)
    return slide

def _timeline_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _heading(slide, data, primary)
    items = (data.get("items") or [])[:5]
    count = max(1, len(items))
    start, usable = .9, 11.5
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(start), Inches(3.35),
           Inches(usable), Inches(.04), secondary)
    step = usable / count
    for i, item in enumerate(items):
        x = start + step*i + step/2
        _shape(slide, MSO_SHAPE.OVAL, Inches(x-.14), Inches(3.21),
               Inches(.28), Inches(.28), primary)
        top = 2.0 if i % 2 == 0 else 3.8
        _text(slide, Inches(x-step/2+.08), Inches(top), Inches(step-.16), Inches(.45),
              item.get("title", ""), 15, INK, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x-step/2+.08), Inches(top+.52), Inches(step-.16), Inches(1.05),
              item.get("body", ""), 11, MUTED, align=PP_ALIGN.CENTER)
    return slide

def _process_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, LIGHT)
    _heading(slide, data, primary)
    items = (data.get("items") or [])[:5]
    for i, item in enumerate(items):
        y = 1.85 + i * .95
        _shape(slide, MSO_SHAPE.OVAL, Inches(.85), Inches(y),
               Inches(.55), Inches(.55), primary)
        _text(slide, Inches(.85), Inches(y+.1), Inches(.55), Inches(.25),
              str(i+1), 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(1.75), Inches(y-.02), Inches(3.3), Inches(.4),
              item.get("title", ""), 17, INK, bold=True, font=HEADER_FONT)
        _text(slide, Inches(5.0), Inches(y-.02), Inches(6.7), Inches(.55),
              item.get("body", ""), 14, MUTED)
    return slide

def _problem_solution_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, WHITE)
    _heading(slide, data, primary)
    _text(slide, Inches(.8), Inches(2.05), Inches(5.2), Inches(.55),
          data.get("left_title", "The Problem"), 22, INK, bold=True, font=HEADER_FONT)
    _bullets(slide, Inches(.8), Inches(2.9), Inches(5.1), Inches(3.1),
             data.get("left_bullets", []), MUTED, 16)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.45), Inches(1.85),
           Inches(5.8), Inches(4.8), primary)
    _text(slide, Inches(6.85), Inches(2.2), Inches(4.9), Inches(.55),
          data.get("right_title", "The Solution"), 22, WHITE, bold=True, font=HEADER_FONT)
    _bullets(slide, Inches(6.85), Inches(3.0), Inches(4.9), Inches(3),
             data.get("right_bullets", []), WHITE, 16)
    return slide

def _three_metrics_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, primary)
    _text(slide, Inches(.75), Inches(.55), Inches(3), Inches(.3),
          str(data.get("eyebrow", "KEY METRICS")).upper(), 10, accent, bold=True)
    _text(slide, Inches(.75), Inches(1.0), Inches(11.5), Inches(.9),
          data.get("heading", ""), 30, WHITE, bold=True, font=HEADER_FONT)
    metrics = (data.get("metrics") or [])[:3]
    for i, metric in enumerate(metrics):
        x = .75 + i*4.1
        _text(slide, Inches(x), Inches(2.7), Inches(3.5), Inches(1.3),
              metric.get("value", ""), 46, WHITE, bold=True, font=HEADER_FONT)
        _shape(slide, MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.15),
               Inches(.8), Inches(.05), accent)
        _text(slide, Inches(x), Inches(4.5), Inches(3.35), Inches(1),
              metric.get("label", ""), 15, WHITE)
    return slide

def _quote_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, primary)
    _text(slide, Inches(.8), Inches(.7), Inches(1), Inches(.8),
          "“", 58, accent, bold=True, font=HEADER_FONT)
    _text(slide, Inches(1.35), Inches(1.65), Inches(10.4), Inches(2.8),
          data.get("quote", ""), 30, WHITE, bold=True, font=HEADER_FONT)
    _text(slide, Inches(1.4), Inches(5.0), Inches(7), Inches(.5),
          f"— {data.get('quote_source', '')}", 14, accent, bold=True)
    return slide

def _conclusion_slide(prs, data, primary, secondary, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, primary)
    _text(slide, Inches(.8), Inches(.7), Inches(3), Inches(.35),
          "KEY TAKEAWAYS", 10, accent, bold=True)
    _text(slide, Inches(.8), Inches(1.35), Inches(10.8), Inches(1.3),
          data.get("heading", "What matters most"), 36, WHITE,
          bold=True, font=HEADER_FONT)
    bullets = (data.get("bullets") or [])[:3]
    for i, item in enumerate(bullets):
        x = .8 + i*4.05
        _text(slide, Inches(x), Inches(3.25), Inches(.55), Inches(.45),
              f"{i+1:02d}", 11, accent, bold=True)
        _text(slide, Inches(x), Inches(3.9), Inches(3.45), Inches(1.6),
              item, 17, WHITE, bold=True, font=HEADER_FONT)
    _text(slide, Inches(.8), Inches(6.55), Inches(4), Inches(.4),
          "THANK YOU", 10, accent, bold=True)
    return slide

def build_presentation(outline: dict, title: str) -> BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    palette = outline.get("palette", {})
    primary = _hex_to_rgb(palette.get("primary"), "#1E2761")
    secondary = _hex_to_rgb(palette.get("secondary"), "#CADCFC")
    accent = _hex_to_rgb(palette.get("accent"), "#F4C95D")

    _title_slide(prs, title, outline.get("subtitle", ""), primary, secondary, accent)

    layout_fns = {
        "hero_statement": _hero_statement_slide,
        "bullets": _bullets_slide,
        "cards": _cards_slide,
        "two_column": _two_column_slide,
        "stat": _stat_slide,
        "timeline": _timeline_slide,
        "process": _process_slide,
        "problem_solution": _problem_solution_slide,
        "three_metrics": _three_metrics_slide,
        "quote": _quote_slide,
    }

    for slide_number, slide_data in enumerate(outline.get("slides", []), start=2):
        layout = slide_data.get("layout", "bullets")
        fn = layout_fns.get(layout, _bullets_slide)
        slide = fn(prs, slide_data, primary, secondary, accent)
        if layout not in {"hero_statement", "three_metrics", "quote"}:
            _footer(slide, slide_number, primary)

    _conclusion_slide(prs, outline.get("conclusion", {}), primary, secondary, accent)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
