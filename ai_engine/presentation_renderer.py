from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_DISPLAY = "Aptos Display"
FONT_BODY = "Aptos"

BACKGROUND = RGBColor(7, 9, 15)
SURFACE = RGBColor(20, 23, 34)
SURFACE_LIGHT = RGBColor(31, 35, 50)

TEXT = RGBColor(245, 246, 248)
MUTED = RGBColor(151, 157, 174)

ACCENT = RGBColor(139, 124, 255)
ACCENT_LIGHT = RGBColor(184, 173, 255)
ACCENT_BLUE = RGBColor(111, 168, 255)


def _set_background(
    slide,
    color=BACKGROUND
) -> None:

    fill = slide.background.fill

    fill.solid()

    fill.fore_color.rgb = color


def _add_text(
    slide,
    *,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color=TEXT,
    bold=False,
    font_name=FONT_BODY,
    align=PP_ALIGN.LEFT
):

    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    frame = box.text_frame

    frame.clear()

    frame.word_wrap = True

    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    paragraph = frame.paragraphs[0]

    paragraph.alignment = align

    run = paragraph.add_run()

    run.text = str(
        text or ""
    )

    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    return box


def _add_shape(
    slide,
    shape_type,
    *,
    left,
    top,
    width,
    height,
    fill_color
):

    shape = slide.shapes.add_shape(
        shape_type,
        left,
        top,
        width,
        height
    )

    shape.fill.solid()

    shape.fill.fore_color.rgb = (
        fill_color
    )

    shape.line.fill.background()

    return shape


def _add_slide_number(
    slide,
    slide_number: int
) -> None:

    _add_text(
        slide,
        left=Inches(11.85),
        top=Inches(6.92),
        width=Inches(0.65),
        height=Inches(0.25),
        text=f"{slide_number:02d}",
        font_size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT
    )


def _add_slide_header(
    slide,
    slide_data: dict,
    slide_number: int
) -> None:

    role = str(
        slide_data.get(
            "role",
            "insight"
        )
    ).replace(
        "_",
        " "
    ).upper()

    _add_text(
        slide,
        left=Inches(0.75),
        top=Inches(0.42),
        width=Inches(3.8),
        height=Inches(0.3),
        text=role,
        font_size=9,
        color=ACCENT_LIGHT,
        bold=True
    )

    _add_text(
        slide,
        left=Inches(0.75),
        top=Inches(0.88),
        width=Inches(11.6),
        height=Inches(0.82),
        text=slide_data.get(
            "title",
            ""
        ),
        font_size=29,
        color=TEXT,
        bold=True,
        font_name=FONT_DISPLAY
    )

    _add_slide_number(
        slide,
        slide_number
    )


def _add_key_message(
    slide,
    slide_data: dict,
    *,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.2),
    height=Inches(1.45)
) -> None:

    _add_text(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        text=slide_data.get(
            "key_message",
            ""
        ),
        font_size=24,
        color=TEXT,
        bold=True,
        font_name=FONT_DISPLAY
    )


def _add_bullets(
    slide,
    bullets: list,
    *,
    left,
    top,
    width,
    height
) -> None:

    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )

    frame = box.text_frame

    frame.clear()

    frame.word_wrap = True

    for index, bullet in enumerate(
        bullets[:4]
    ):

        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )

        paragraph.text = (
            f"•  {bullet}"
        )

        paragraph.font.name = (
            FONT_BODY
        )

        paragraph.font.size = Pt(15)

        paragraph.font.color.rgb = (
            TEXT
        )

        paragraph.space_after = Pt(13)


def _render_typography(
    slide,
    slide_data: dict
) -> None:

    subtitle = slide_data.get(
        "subtitle",
        ""
    )

    _add_text(
        slide,
        left=Inches(0.9),
        top=Inches(1.55),
        width=Inches(10.8),
        height=Inches(2.1),
        text=slide_data.get(
            "title",
            ""
        ),
        font_size=40,
        color=TEXT,
        bold=True,
        font_name=FONT_DISPLAY
    )

    _add_text(
        slide,
        left=Inches(0.92),
        top=Inches(4.0),
        width=Inches(8.8),
        height=Inches(1.0),
        text=(
            subtitle
            or slide_data.get(
                "key_message",
                ""
            )
        ),
        font_size=18,
        color=MUTED
    )

    _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left=Inches(0.92),
        top=Inches(5.55),
        width=Inches(1.0),
        height=Inches(0.06),
        fill_color=ACCENT
    )


def _render_structured_text(
    slide,
    slide_data: dict
) -> None:

    _add_key_message(
        slide,
        slide_data
    )

    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(6.45),
        top=Inches(1.95),
        width=Inches(5.8),
        height=Inches(4.4),
        fill_color=SURFACE
    )

    _add_bullets(
        slide,
        slide_data.get(
            "bullets",
            []
        ),
        left=Inches(6.9),
        top=Inches(2.45),
        width=Inches(4.9),
        height=Inches(3.3)
    )


def _render_process_flow(
    slide,
    slide_data: dict,
    visual_data: dict
) -> None:

    nodes = visual_data.get(
        "diagram_nodes",
        []
    )

    if not nodes:
        _render_structured_text(
            slide,
            slide_data
        )

        return

    _add_text(
        slide,
        left=Inches(0.8),
        top=Inches(1.72),
        width=Inches(11.5),
        height=Inches(0.7),
        text=slide_data.get(
            "key_message",
            ""
        ),
        font_size=19,
        color=MUTED
    )

    count = len(nodes)

    usable_width = 11.4

    node_width = min(
        2.25,
        (
            usable_width
            - (count - 1) * 0.5
        ) / count
    )

    total_width = (
        node_width * count
        + 0.5 * (count - 1)
    )

    start_x = (
        13.333 - total_width
    ) / 2

    for index, node in enumerate(nodes):

        x = (
            start_x
            + index * (
                node_width + 0.5
            )
        )

        if index < count - 1:

            connector = (
                slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(
                        x + node_width
                    ),
                    Inches(4.05),
                    Inches(
                        x
                        + node_width
                        + 0.5
                    ),
                    Inches(4.05)
                )
            )

            connector.line.color.rgb = (
                ACCENT
            )

            connector.line.width = Pt(1.5)

        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(x),
            top=Inches(3.25),
            width=Inches(node_width),
            height=Inches(1.65),
            fill_color=SURFACE_LIGHT
        )

        _add_text(
            slide,
            left=Inches(x + 0.18),
            top=Inches(3.58),
            width=Inches(
                node_width - 0.36
            ),
            height=Inches(0.8),
            text=node,
            font_size=15,
            color=TEXT,
            bold=True,
            align=PP_ALIGN.CENTER
        )


def _render_comparison(
    slide,
    slide_data: dict
) -> None:

    bullets = slide_data.get(
        "bullets",
        []
    )

    midpoint = max(
        1,
        len(bullets) // 2
    )

    left_bullets = bullets[
        :midpoint
    ]

    right_bullets = bullets[
        midpoint:
    ]

    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(0.75),
        top=Inches(1.95),
        width=Inches(5.75),
        height=Inches(4.55),
        fill_color=SURFACE
    )

    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left=Inches(6.8),
        top=Inches(1.95),
        width=Inches(5.75),
        height=Inches(4.55),
        fill_color=SURFACE_LIGHT
    )

    _add_text(
        slide,
        left=Inches(1.1),
        top=Inches(2.3),
        width=Inches(4.9),
        height=Inches(0.8),
        text="Perspective A",
        font_size=20,
        color=ACCENT_LIGHT,
        bold=True,
        font_name=FONT_DISPLAY
    )

    _add_text(
        slide,
        left=Inches(7.15),
        top=Inches(2.3),
        width=Inches(4.9),
        height=Inches(0.8),
        text="Perspective B",
        font_size=20,
        color=ACCENT_BLUE,
        bold=True,
        font_name=FONT_DISPLAY
    )

    _add_bullets(
        slide,
        left_bullets,
        left=Inches(1.1),
        top=Inches(3.15),
        width=Inches(4.8),
        height=Inches(2.7)
    )

    _add_bullets(
        slide,
        right_bullets,
        left=Inches(7.15),
        top=Inches(3.15),
        width=Inches(4.8),
        height=Inches(2.7)
    )


def _render_diagram(
    slide,
    slide_data: dict,
    visual_data: dict
) -> None:

    nodes = visual_data.get(
        "diagram_nodes",
        []
    )

    if not nodes:
        _render_structured_text(
            slide,
            slide_data
        )

        return

    center_x = 8.9
    center_y = 4.05

    _add_key_message(
        slide,
        slide_data,
        width=Inches(4.7)
    )

    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        left=Inches(center_x - 0.8),
        top=Inches(center_y - 0.8),
        width=Inches(1.6),
        height=Inches(1.6),
        fill_color=ACCENT
    )

    _add_text(
        slide,
        left=Inches(center_x - 0.65),
        top=Inches(center_y - 0.3),
        width=Inches(1.3),
        height=Inches(0.6),
        text=slide_data.get(
            "title",
            ""
        ),
        font_size=12,
        color=TEXT,
        bold=True,
        align=PP_ALIGN.CENTER
    )

    positions = [
        (6.35, 2.25),
        (10.25, 2.25),
        (6.35, 5.0),
        (10.25, 5.0),
    ]

    for index, node in enumerate(
        nodes[:4]
    ):

        x, y = positions[index]

        connector = (
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(center_x),
                Inches(center_y),
                Inches(x + 0.75),
                Inches(y + 0.45)
            )
        )

        connector.line.color.rgb = (
            ACCENT_LIGHT
        )

        connector.line.width = Pt(1.2)

        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(x),
            top=Inches(y),
            width=Inches(1.5),
            height=Inches(0.9),
            fill_color=SURFACE_LIGHT
        )

        _add_text(
            slide,
            left=Inches(x + 0.12),
            top=Inches(y + 0.23),
            width=Inches(1.26),
            height=Inches(0.4),
            text=node,
            font_size=11,
            color=TEXT,
            bold=True,
            align=PP_ALIGN.CENTER
        )


def _render_framework(
    slide,
    slide_data: dict,
    visual_data: dict
) -> None:

    nodes = visual_data.get(
        "diagram_nodes",
        []
    )

    if not nodes:

        nodes = slide_data.get(
            "bullets",
            []
        )

    _add_text(
        slide,
        left=Inches(0.8),
        top=Inches(1.72),
        width=Inches(11.5),
        height=Inches(0.7),
        text=slide_data.get(
            "key_message",
            ""
        ),
        font_size=18,
        color=MUTED
    )

    nodes = nodes[:4]

    count = max(
        1,
        len(nodes)
    )

    card_width = (
        11.5
        - 0.3 * (
            count - 1
        )
    ) / count

    for index, node in enumerate(nodes):

        x = (
            0.8
            + index * (
                card_width + 0.3
            )
        )

        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(x),
            top=Inches(3.0),
            width=Inches(card_width),
            height=Inches(2.4),
            fill_color=(
                SURFACE_LIGHT
                if index % 2
                else SURFACE
            )
        )

        _add_text(
            slide,
            left=Inches(x + 0.25),
            top=Inches(3.4),
            width=Inches(
                card_width - 0.5
            ),
            height=Inches(0.35),
            text=f"{index + 1:02d}",
            font_size=10,
            color=ACCENT_LIGHT,
            bold=True
        )

        _add_text(
            slide,
            left=Inches(x + 0.25),
            top=Inches(4.05),
            width=Inches(
                card_width - 0.5
            ),
            height=Inches(0.8),
            text=node,
            font_size=16,
            color=TEXT,
            bold=True,
            font_name=FONT_DISPLAY
        )


def _render_slide(
    presentation,
    slide_data: dict,
    visual_data: dict
) -> None:

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    _set_background(
        slide
    )

    slide_number = slide_data[
        "slide_number"
    ]

    visual_type = visual_data.get(
        "visual_type",
        "structured_text"
    )

    if visual_type == "typography":

        _render_typography(
            slide,
            slide_data
        )

        _add_slide_number(
            slide,
            slide_number
        )

        return

    _add_slide_header(
        slide,
        slide_data,
        slide_number
    )

    if visual_type == "structured_text":

        _render_structured_text(
            slide,
            slide_data
        )

    elif visual_type == "process_flow":

        _render_process_flow(
            slide,
            slide_data,
            visual_data
        )

    elif visual_type == "comparison":

        _render_comparison(
            slide,
            slide_data
        )

    elif visual_type in {
        "hierarchy",
        "relationship_diagram",
        "conceptual_diagram",
    }:

        _render_diagram(
            slide,
            slide_data,
            visual_data
        )

    elif visual_type in {
        "synthesis_framework",
        "timeline",
    }:

        _render_framework(
            slide,
            slide_data,
            visual_data
        )

    else:

        _render_structured_text(
            slide,
            slide_data
        )


def render_presentation(
    generation_package: dict
) -> BytesIO:

    if not isinstance(
        generation_package,
        dict
    ):
        raise TypeError(
            "generation_package must be a dictionary."
        )

    slide_deck = generation_package.get(
        "slide_deck"
    )

    visual_plan = generation_package.get(
        "visual_plan"
    )

    if not isinstance(
        slide_deck,
        dict
    ):
        raise ValueError(
            "Generation package is missing slide_deck."
        )

    if not isinstance(
        visual_plan,
        dict
    ):
        raise ValueError(
            "Generation package is missing visual_plan."
        )

    content_slides = slide_deck.get(
        "slides"
    )

    visual_slides = visual_plan.get(
        "slides"
    )

    if not isinstance(
        content_slides,
        list
    ):
        raise ValueError(
            "slide_deck slides must be a list."
        )

    if not isinstance(
        visual_slides,
        list
    ):
        raise ValueError(
            "visual_plan slides must be a list."
        )

    if len(content_slides) != len(
        visual_slides
    ):
        raise ValueError(
            "Slide content and visual plan counts do not match."
        )

    visual_by_number = {
        slide["slide_number"]: slide
        for slide in visual_slides
    }

    presentation = Presentation()

    presentation.slide_width = (
        SLIDE_WIDTH
    )

    presentation.slide_height = (
        SLIDE_HEIGHT
    )

    for slide_data in content_slides:

        slide_number = slide_data.get(
            "slide_number"
        )

        visual_data = visual_by_number.get(
            slide_number
        )

        if visual_data is None:
            raise ValueError(
                "Missing visual plan for "
                f"slide {slide_number}."
            )

        _render_slide(
            presentation,
            slide_data,
            visual_data
        )

    buffer = BytesIO()

    presentation.save(
        buffer
    )

    buffer.seek(
        0
    )

    print(
        "[PRESENTATION RENDERER] "
        f"Rendered {len(content_slides)} slides."
    )

    return buffer